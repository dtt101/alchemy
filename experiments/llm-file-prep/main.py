import mammoth
from pptx import Presentation
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
from pathlib import Path

HERE = Path(__file__).resolve().parent
DOCX = HERE / "test.docx"
PPTX = HERE / "test.pptx"
PDFDOC = HERE / "test.pdf"
with open(DOCX, "rb") as test_doc:
    result = mammoth.convert_to_html(test_doc)
    html = result.value
    messages = result.messages

print('MAMMOTH')
print(html)
print(messages)

# TODO 
# python-pptx for ppt
text_runs = []
prs = Presentation(PPTX)
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

print("PYTHON-PPTX")
print(text_runs)

# pdfminer.six for pdf
print("PDFMINER")
for page_layout in extract_pages(PDFDOC):
    for element in page_layout:
        if isinstance(element, LTTextContainer):
            print(element.get_text())


# TODO - convert HTML to markdown with beautifulsoup
